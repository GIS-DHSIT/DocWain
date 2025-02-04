import fitz  # PyMuPDF for PDFs
import pandas as pd
from pptx import Presentation
import docx

def extract_text(file_path):
    if file_path.endswith(".pdf"):
        doc = fitz.open(file_path)
        return "\n".join([page.get_text() for page in doc])

    elif file_path.endswith(".docx"):
        doc = docx.Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])

    elif file_path.endswith(".csv") or file_path.endswith(".xlsx"):
        df = pd.read_csv(file_path) if file_path.endswith(".csv") else pd.read_excel(file_path)
        return df.to_string()

    elif file_path.endswith(".pptx"):
        prs = Presentation(file_path)
        return "\n".join([slide.shapes.title.text for slide in prs.slides if slide.shapes.title])

    return ""
