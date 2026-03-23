import hashlib
import io
from src.utils.logging_utils import get_logger
import re
from typing import Any, Dict, List, Optional, Tuple, Union

import docx
try:
    import pymupdf as fitz  # PyMuPDF >= 1.24 exports as pymupdf
except ImportError:
    import fitz  # Older PyMuPDF versions export as fitz
import numpy as np
import pandas as pd
import pytesseract
from PIL import Image
from pptx import Presentation
try:
    import pdfplumber
except Exception:  # noqa: BLE001
    pdfplumber = None
from src.api.config import Config
from src.api.pipeline_models import ChunkCandidate, ExtractedDocument, Figure, Section, Table

logger = get_logger(__name__)

try:
    import easyocr
except Exception:  # noqa: BLE001
    easyocr = None

_OCR_RETRY_THRESHOLD = 70.0
_DIGIT_MASK_RE = re.compile(r"\d+")

def _smart_decode(raw: bytes) -> str:
    """Decode bytes to str with intelligent encoding detection.

    Tries: BOM detection -> null-ratio heuristic -> charset_normalizer -> UTF-8 -> Latin-1.
    """
    # 1. BOM detection (UTF-16LE, UTF-16BE, UTF-8-sig)
    if raw[:2] == b'\xff\xfe':
        return raw.decode('utf-16-le', errors='replace')
    if raw[:2] == b'\xfe\xff':
        return raw.decode('utf-16-be', errors='replace')
    if raw[:3] == b'\xef\xbb\xbf':
        return raw.decode('utf-8-sig', errors='replace')

    # 2. Check for UTF-16LE without BOM (common in .doc files):
    #    ASCII chars encoded as char+\x00 -> every other byte is \x00
    if len(raw) >= 20:
        sample = raw[:100]
        null_ratio = sample.count(b'\x00') / len(sample)
        if null_ratio > 0.3:
            try:
                return raw.decode('utf-16', errors='replace')
            except Exception:
                pass

    # 3. charset_normalizer auto-detection
    try:
        from charset_normalizer import from_bytes
        result = from_bytes(raw).best()
        if result is not None and result.encoding:
            return str(result)
    except Exception:
        pass

    # 4. UTF-8 fallback
    try:
        return raw.decode('utf-8', errors='replace')
    except Exception:
        return raw.decode('latin-1', errors='replace')
_COPYRIGHT_RE = re.compile(r"(?:copyright|©|\(c\))\s*\d{4}", re.IGNORECASE)
_CONFIDENTIAL_RE = re.compile(r"\b(?:confidential|proprietary|internal use only|do not distribute)\b", re.IGNORECASE)

class DocumentExtractor:
    """Layout-aware document extractor with selective OCR and structured output."""

    def __init__(self, ocr_engine: Optional[str] = None):
        self.ocr_engine = (ocr_engine or getattr(Config.Model, "OCR_ENGINE", "pytesseract")).lower()
        self._easyocr_reader = None
        self._doc_intel = DocumentIntelligence()

    # ---------- Helpers ----------
    @staticmethod
    def _is_heading(text: str) -> bool:
        if not text:
            return False
        clean = text.strip()
        if len(clean.split()) <= 1 and clean.isupper():
            return True
        heading_pattern = re.compile(r"^(\d+\.)+\s+.+|^[A-Z][A-Z0-9\s,:-]{4,}$")
        return bool(heading_pattern.match(clean))

    @staticmethod
    def _looks_like_table(text: str) -> bool:
        if not text:
            return False
        separators = text.count("|") + text.count(",") + text.count("\t")
        return separators >= 4

    @staticmethod
    def _make_section_id(title: str, page: int) -> str:
        raw = f"{title}|{page}"
        return hashlib.sha1(raw.encode("utf-8")).hexdigest()[:12]

    @staticmethod
    def _detect_scanned_page(text: str, image_count: int) -> bool:
        return (not text or len(text.strip()) < 30) and image_count > 0

    @staticmethod
    def _page_number_line(text: str) -> bool:
        return bool(re.match(r"^\s*(?:page\s*)?\d+(?:\s*(?:/|of)\s*\d+)?\s*$", text, re.IGNORECASE))

    @staticmethod
    def _extract_key_value_pairs(text: str) -> List[Dict[str, Any]]:
        pairs: List[Dict[str, Any]] = []
        if not text:
            return pairs
        for line in text.splitlines():
            candidate = line.strip()
            if not candidate or len(candidate) > 200:
                continue
            match = re.match(r"^([A-Za-z0-9][^:]{1,64}):\s*(.+)$", candidate)
            if not match:
                match = re.match(r"^([A-Za-z0-9][^-]{1,64})\s+-\s+(.+)$", candidate)
            if match:
                key = match.group(1).strip()
                value = match.group(2).strip()
                if key and value:
                    pairs.append({"key": key, "value": value})
        return pairs

    @staticmethod
    def _layout_confidence(chunks: List[ChunkCandidate]) -> float:
        if not chunks:
            return 0.0
        structured = sum(1 for cand in chunks if cand.chunk_type in {"table", "table_row", "table_header"})
        density = structured / max(1, len(chunks))
        return round(min(1.0, 0.6 + density * 0.4), 3)

    def _assess_doc_quality(
        self,
        *,
        ocr_confidences: List[float],
        chunk_candidates: List[ChunkCandidate],
        has_images: bool,
    ) -> Dict[str, Any]:
        min_conf = float(getattr(Config.Retrieval, "MIN_OCR_CONFIDENCE", 60))
        avg_conf = None
        if ocr_confidences:
            avg_conf = sum(ocr_confidences) / max(1, len(ocr_confidences))
        layout_conf = self._layout_confidence(chunk_candidates)
        if avg_conf is None:
            quality = "MEDIUM" if chunk_candidates else "LOW"
        elif avg_conf < min_conf:
            quality = "LOW"
        elif avg_conf < min_conf + 15:
            quality = "MEDIUM"
        else:
            quality = "HIGH"
        return {
            "doc_quality": quality,
            "ocr_confidence_avg": avg_conf,
            "layout_confidence": layout_conf,
            "ocr_upgrade_suggested": bool(has_images and quality == "LOW"),
        }

    def _ensure_easyocr(self):
        if easyocr and self._easyocr_reader is None:
            self._easyocr_reader = easyocr.Reader(["en"], gpu=False)

    def _ocr_pytesseract(self, image: Image.Image) -> Tuple[str, Optional[float]]:
        """Run pytesseract OCR on *image* and return (text, confidence%)."""
        try:
            data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
            words = [w.strip() for w in data.get("text", []) if w.strip()]
            confidences = []
            for conf in data.get("conf", []):
                try:
                    conf_val = float(conf)
                except (TypeError, ValueError):
                    continue
                if conf_val >= 0:
                    confidences.append(conf_val)
            text_value = " ".join(words).strip()
            avg_conf = sum(confidences) / len(confidences) if confidences else None
            return text_value, avg_conf
        except Exception as exc:  # noqa: BLE001
            logger.error("pytesseract OCR failed: %s", exc)
            return "", None

    def _ocr_easyocr(self, image: Image.Image) -> Tuple[str, Optional[float]]:
        """Run easyocr on *image* and return (text, confidence%)."""
        try:
            self._ensure_easyocr()
            result = self._easyocr_reader.readtext(np.array(image), detail=1)
            texts = []
            confidences = []
            for _, text, conf in result:
                if text:
                    texts.append(text)
                if conf is not None:
                    confidences.append(float(conf) * 100.0)
            text_value = " ".join(texts).strip()
            avg_conf = sum(confidences) / len(confidences) if confidences else None
            return text_value, avg_conf
        except Exception as exc:  # noqa: BLE001
            logger.warning("EasyOCR failed: %s", exc)
            return "", None

    def _ocr_image(self, image: Image.Image, *, engine: Optional[str] = None) -> Tuple[str, Optional[float]]:
        """Coordinate OCR with per-image dual-engine retry.

        Calls the primary engine first (default: pytesseract).  If the
        confidence is below ``_OCR_RETRY_THRESHOLD`` (70%), retries with
        the alternate engine and returns whichever result is better.
        """
        active_engine = (engine or self.ocr_engine).lower()

        # --- primary pass ---
        if active_engine == "easyocr" and easyocr:
            primary_text, primary_conf = self._ocr_easyocr(image)
            alternate_method = self._ocr_pytesseract
        else:
            primary_text, primary_conf = self._ocr_pytesseract(image)
            alternate_method = self._ocr_easyocr if easyocr else None

        # --- per-image retry when confidence is low ---
        if (
            alternate_method is not None
            and (primary_conf is None or primary_conf < _OCR_RETRY_THRESHOLD)
        ):
            try:
                retry_text, retry_conf = alternate_method(image)
                # Pick the better result (higher confidence wins; treat None as 0)
                retry_conf_val = retry_conf if retry_conf is not None else 0.0
                primary_conf_val = primary_conf if primary_conf is not None else 0.0
                if retry_conf_val > primary_conf_val:
                    return retry_text, retry_conf
            except Exception as exc:  # noqa: BLE001
                logger.debug("OCR retry with alternate engine failed: %s", exc)

        return primary_text, primary_conf

    def _ocr_with_vision(
        self, image: Image.Image, *, is_full_page: bool = False,
    ) -> Tuple[str, Optional[float]]:
        """Run vision OCR (glm-ocr) with traditional OCR fallback.

        When both produce results, picks the best via ``_pick_best_ocr``.
        Falls through to ``_ocr_image()`` when vision is unavailable.
        """
        from src.api.config import Config

        cfg = getattr(Config, "VisionOCR", None)
        if not cfg or not getattr(cfg, "ENABLED", True):
            return self._ocr_image(image)

        try:
            from src.llm.vision_ocr import get_vision_ocr_client

            client = get_vision_ocr_client()
            if client is None or not client.is_available():
                return self._ocr_image(image)

            if is_full_page:
                vision_text, vision_conf = client.ocr_page_image(image)
            else:
                vision_text, vision_conf = client.ocr_image(image)
        except Exception as exc:
            logger.debug("Vision OCR unavailable, falling back to traditional: %s", exc)
            return self._ocr_image(image)

        # Also run traditional OCR for comparison / fallback
        if getattr(cfg, "FALLBACK_TO_TRADITIONAL", True):
            trad_text, trad_conf = self._ocr_image(image)
            return self._pick_best_ocr(vision_text, vision_conf, trad_text, trad_conf)

        if vision_text:
            return vision_text, vision_conf
        # Vision produced nothing — fall back
        return self._ocr_image(image)

    @staticmethod
    def _pick_best_ocr(
        vision_text: str,
        vision_conf: Optional[float],
        trad_text: str,
        trad_conf: Optional[float],
    ) -> Tuple[str, Optional[float]]:
        """Compare vision vs traditional OCR results and return the best."""
        v_len = len(vision_text.strip()) if vision_text else 0
        t_len = len(trad_text.strip()) if trad_text else 0

        # One empty, other not → pick non-empty
        if v_len and not t_len:
            return vision_text, vision_conf
        if t_len and not v_len:
            return trad_text, trad_conf
        if not v_len and not t_len:
            return "", None

        # Vision captured >50% more text → prefer vision
        if v_len > t_len * 1.5:
            return vision_text, vision_conf

        # Traditional much longer AND high confidence → prefer traditional
        trad_conf_val = trad_conf if trad_conf is not None else 0.0
        if t_len > v_len * 1.5 and trad_conf_val >= 80.0:
            return trad_text, trad_conf

        # Default: prefer vision (better contextual understanding)
        return vision_text, vision_conf

    def _finalize_section(
        self,
        sections: List[Section],
        title: str,
        section_id: str,
        start_page: int,
        end_page: int,
        buffer: List[str],
    ):
        if buffer:
            sections.append(
                Section(
                    section_id=section_id,
                    title=title,
                    level=1,
                    start_page=start_page,
                    end_page=end_page,
                    text="\n".join(buffer).strip(),
                )
            )

    @staticmethod
    def _reorder_text_blocks(blocks: list, page_width: float) -> list:
        """Reorder raw PyMuPDF block tuples for correct multi-column reading order.

        PyMuPDF blocks are tuples: (x0, y0, x1, y1, text, block_no, block_type).
        Multi-column PDFs (resumes, academic papers) return blocks interleaved
        across columns. This method detects columns and reads each top-to-bottom
        before moving to the next.

        Strategy:
        1. Separate text blocks (type 0) from image blocks (type 1)
        2. Detect columns by clustering block x-midpoints
        3. Full-width blocks (>55% page width) read first in y-order
        4. Left column blocks read top-to-bottom
        5. Right column blocks read top-to-bottom
        6. Supports 2-column and 3-column layouts
        """
        if not blocks or page_width <= 0:
            return blocks

        # Filter to text blocks only (type 0)
        text_blocks = [b for b in blocks if len(b) >= 7 and b[6] == 0]
        other_blocks = [b for b in blocks if len(b) < 7 or b[6] != 0]

        if len(text_blocks) < 3:
            return blocks  # Too few to rearrange

        mid = page_width / 2.0
        gap_tolerance = page_width * 0.05  # 5% of page width

        # Classify blocks by position
        full_width = []
        left_col = []
        right_col = []

        for b in text_blocks:
            x0, y0, x1, y1 = b[0], b[1], b[2], b[3]
            block_width = x1 - x0
            block_mid = (x0 + x1) / 2.0

            if block_width > page_width * 0.55:
                full_width.append(b)
            elif block_mid < mid + gap_tolerance:
                left_col.append(b)
            else:
                right_col.append(b)

        # Check if this is actually multi-column
        is_multi = len(left_col) >= 2 and len(right_col) >= 2

        if not is_multi:
            # Single column — just sort by vertical position
            return sorted(text_blocks, key=lambda b: (b[1], b[0])) + other_blocks

        # Multi-column: full-width headers first, then left column, then right column
        full_width.sort(key=lambda b: b[1])
        left_col.sort(key=lambda b: b[1])
        right_col.sort(key=lambda b: b[1])

        # Interleave full-width blocks at their correct vertical position
        result = []
        fw_idx = 0
        l_idx = 0
        r_idx = 0

        # First pass: full-width blocks that appear above both columns
        while fw_idx < len(full_width):
            fw_y = full_width[fw_idx][1]
            first_left_y = left_col[0][1] if left_col else float('inf')
            first_right_y = right_col[0][1] if right_col else float('inf')
            if fw_y < min(first_left_y, first_right_y):
                result.append(full_width[fw_idx])
                fw_idx += 1
            else:
                break

        # Left column
        result.extend(left_col)

        # Right column
        result.extend(right_col)

        # Remaining full-width blocks (footers)
        result.extend(full_width[fw_idx:])

        # Add back non-text blocks
        result.extend(other_blocks)

        return result

    def _reorder_blocks_by_layout(self, blocks: List[dict], page_width: float = 612.0) -> List[dict]:
        """Reorder text blocks using layout analysis for multi-column detection."""
        try:
            from src.intelligence.layout_analyzer import LayoutAnalyzer
            analyzer = LayoutAnalyzer(page_width=page_width)
            result = analyzer.analyze(blocks)
            if result.ordered_blocks:
                return result.ordered_blocks
        except Exception:
            pass
        return blocks

    @staticmethod
    def _coerce_file_like(content: object):
        """
        Normalize raw payloads into a seekable stream for libraries like python-docx/pptx.
        Handles bytes, memoryview, and already-open file-like objects.
        """
        if isinstance(content, memoryview):
            content = content.tobytes()
        if isinstance(content, (bytes, bytearray)):
            return io.BytesIO(content)
        if isinstance(content, str):
            # Treat as raw text payload, not a path; decode to bytes
            return io.BytesIO(content.encode("utf-8"))
        if hasattr(content, "read"):
            try:
                content.seek(0)
            except Exception:
                pass
            return content
        try:
            return io.BytesIO(bytes(content))
        except Exception:
            return io.BytesIO()

    def _build_canonical_json(
        self,
        *,
        pages: List[Dict[str, Any]],
        sections: List[Section],
        tables: List[Table],
        figures: List[Figure],
        chunk_candidates: List[ChunkCandidate],
        layout_blocks: Optional[List[Dict[str, Any]]] = None,
        page_dims: Optional[Dict[int, Dict[str, float]]] = None,
    ) -> Dict[str, Any]:
        key_values: List[Dict[str, Any]] = []
        for section in sections:
            for pair in self._extract_key_value_pairs(section.text):
                key_values.append(
                    {
                        "key": pair["key"],
                        "value": pair["value"],
                        "section_title": section.title,
                        "page_start": section.start_page,
                        "page_end": section.end_page,
                    }
                )
        layout_spans = [
            {
                "section_id": cand.section_id,
                "section_title": cand.section_title,
                "page": cand.page,
                "chunk_type": cand.chunk_type,
                "text": cand.text,
            }
            for cand in chunk_candidates
            if cand.text
        ]
        payload: Dict[str, Any] = {
            "pages": pages,
            "sections": [
                {
                    "section_id": sec.section_id,
                    "title": sec.title,
                    "level": sec.level,
                    "start_page": sec.start_page,
                    "end_page": sec.end_page,
                    "text": sec.text,
                }
                for sec in sections
            ],
            "tables": [
                {"page": tbl.page, "text": tbl.text, "csv": tbl.csv} for tbl in tables
            ],
            "figures": [
                {"page": fig.page, "caption": fig.caption} for fig in figures
            ],
            "key_value_pairs": key_values,
            "layout_spans": layout_spans,
        }
        if layout_blocks:
            payload["layout_blocks"] = layout_blocks
        if page_dims:
            payload["page_dims"] = page_dims
        return payload

    @staticmethod
    def _dedupe_page_lines(pages: Dict[int, List[str]]) -> Tuple[Dict[int, List[str]], List[str], List[str]]:
        if not pages:
            return {}, [], []
        headers: Dict[str, int] = {}
        footers: Dict[str, int] = {}
        # Also track digit-masked versions for running headers/footers with page numbers
        masked_headers: Dict[str, int] = {}
        masked_footers: Dict[str, int] = {}

        for lines in pages.values():
            non_empty = [ln.strip() for ln in lines if ln.strip()]
            if not non_empty:
                continue
            first = non_empty[0]
            last = non_empty[-1]

            headers[first] = headers.get(first, 0) + 1
            footers[last] = footers.get(last, 0) + 1

            # Mask digits to group "Page 1", "Page 2", etc. as same pattern
            masked_first = _DIGIT_MASK_RE.sub("N", first)
            masked_last = _DIGIT_MASK_RE.sub("N", last)
            masked_headers[masked_first] = masked_headers.get(masked_first, 0) + 1
            masked_footers[masked_last] = masked_footers.get(masked_last, 0) + 1

        total_pages = max(1, len(pages))
        threshold = 0.6
        header_lines_set = {line for line, count in headers.items() if count / total_pages >= threshold}
        footer_lines_set = {line for line, count in footers.items() if count / total_pages >= threshold}

        # Running headers/footers (digit-masked patterns)
        masked_header_patterns = {pat for pat, count in masked_headers.items() if count / total_pages >= threshold}
        masked_footer_patterns = {pat for pat, count in masked_footers.items() if count / total_pages >= threshold}

        # Build the return lists (for backward compat)
        header_lines = list(header_lines_set)
        footer_lines = list(footer_lines_set)

        cleaned: Dict[int, List[str]] = {}
        for page, lines in pages.items():
            filtered = []
            non_empty = [ln.strip() for ln in lines if ln.strip()]
            first_line = non_empty[0] if non_empty else ""
            last_line = non_empty[-1] if non_empty else ""

            for line in lines:
                stripped = line.strip()
                if not stripped:
                    continue
                # Exact match headers/footers
                if stripped in header_lines_set or stripped in footer_lines_set:
                    continue
                # Page number lines
                if DocumentExtractor._page_number_line(stripped):
                    continue
                # Digit-masked pattern match (for running headers with page numbers)
                masked = _DIGIT_MASK_RE.sub("N", stripped)
                if stripped == first_line and masked in masked_header_patterns:
                    continue
                if stripped == last_line and masked in masked_footer_patterns:
                    continue
                # Copyright lines
                if _COPYRIGHT_RE.search(stripped):
                    continue
                # Confidentiality notices (only short ones, to avoid stripping content)
                if _CONFIDENTIAL_RE.search(stripped) and len(stripped) < 60:
                    continue
                filtered.append(line)
            cleaned[page] = filtered
        return cleaned, header_lines, footer_lines

    # ---------- Extraction routines ----------
    def extract_text_from_pdf(self, pdf_content: bytes, filename: Optional[str] = None) -> ExtractedDocument:
        sections: List[Section] = []
        tables: List[Table] = []
        figures: List[Figure] = []
        chunk_candidates: List[ChunkCandidate] = []
        full_text_parts: List[str] = []
        errors: List[str] = []
        ocr_confidences: List[float] = []
        pages_map: Dict[int, List[str]] = {}
        ocr_targets: List[Dict[str, Any]] = []
        layout_blocks: List[Dict[str, Any]] = []
        page_dims: Dict[int, Dict[str, float]] = {}

        current_title = "Introduction"
        current_section_id = self._make_section_id(current_title, 1)
        current_start_page = 1
        section_buffer: List[str] = []
        last_page = 1

        if hasattr(fitz, "open"):
            try:
                with fitz.open(stream=pdf_content, filetype="pdf") as doc:
                    for page_index, page in enumerate(doc, start=1):
                        last_page = page_index
                        try:
                            page_dims[page_index] = {"width": float(page.rect.width), "height": float(page.rect.height)}
                        except Exception:
                            page_dims[page_index] = {}
                        blocks = page.get_text("blocks", sort=True) or []
                        page_text_parts: List[str] = []

                        # ── Multi-column layout detection & reordering ──
                        # Sort blocks spatially for proper reading order in multi-column PDFs
                        page_width = float(page.rect.width)
                        if len(blocks) >= 4 and page_width > 0:
                            blocks = self._reorder_text_blocks(blocks, page_width)

                        for block in blocks:
                            if len(block) < 5:
                                continue
                            block_text = (block[4] or "").strip()
                            if not block_text:
                                continue

                            first_line = block_text.splitlines()[0].strip()
                            if self._is_heading(first_line):
                                self._finalize_section(
                                    sections, current_title, current_section_id, current_start_page, page_index, section_buffer
                                )
                                section_buffer = []
                                current_title = first_line
                                current_section_id = self._make_section_id(first_line, page_index)
                                current_start_page = page_index
                                continue

                            section_buffer.append(block_text)
                            chunk_type = "table" if self._looks_like_table(block_text) else "text"
                            chunk_candidates.append(
                                ChunkCandidate(
                                    text=block_text,
                                    page=page_index,
                                    section_title=current_title,
                                    section_id=current_section_id,
                                    chunk_type=chunk_type,
                                )
                            )
                            page_text_parts.append(block_text)
                            if chunk_type == "table":
                                tables.append(Table(page=page_index, text=block_text))

                        # Capture layout blocks with geometry + typography when available.
                        try:
                            layout_dict = page.get_text("dict")
                            for block in layout_dict.get("blocks", []) or []:
                                block_bbox = block.get("bbox")
                                block_type = "image" if block.get("type") == 1 else "text"
                                text_parts: List[str] = []
                                font_sizes: List[float] = []
                                bold = False
                                line_count = 0
                                if block_type == "text":
                                    for line in block.get("lines", []) or []:
                                        line_count += 1
                                        for span in line.get("spans", []) or []:
                                            span_text = str(span.get("text") or "")
                                            if span_text:
                                                text_parts.append(span_text)
                                            size_val = span.get("size")
                                            if isinstance(size_val, (int, float)):
                                                font_sizes.append(float(size_val))
                                            font_name = str(span.get("font") or "").lower()
                                            if "bold" in font_name or "black" in font_name:
                                                bold = True
                                text_value = " ".join(text_parts).strip()
                                if not text_value and block_type != "image":
                                    continue
                                first_line = text_value.splitlines()[0].strip() if text_value else ""
                                layout_blocks.append(
                                    {
                                        "page": page_index,
                                        "bbox": block_bbox or [0, 0, 0, 0],
                                        "text": text_value,
                                        "block_type": "table"
                                        if self._looks_like_table(first_line)
                                        else ("list" if re.match(r"^\s*(?:[-*•]|\d+[\.)])\s+", first_line) else block_type),
                                        "style": {
                                            "font_size": (sum(font_sizes) / len(font_sizes)) if font_sizes else None,
                                            "font_weight": "bold" if bold else "normal",
                                            "is_all_caps": bool(text_value) and text_value.isupper(),
                                            "line_count": line_count or None,
                                        },
                                    }
                                )
                        except Exception as exc:  # noqa: BLE001
                            errors.append(f"pdf_layout_blocks_failed: page={page_index} err={exc}")

                        # Apply layout analysis to reorder multi-column content
                        try:
                            page_layout_blocks = [b for b in layout_blocks if b.get("page") == page_index]
                            if page_layout_blocks:
                                reordered = self._reorder_blocks_by_layout(page_layout_blocks, page_width=float(page.rect.width))
                                # Replace the page's layout blocks with reordered ones
                                layout_blocks = [b for b in layout_blocks if b.get("page") != page_index] + reordered
                        except Exception:
                            pass

                        images = page.get_images(full=True)
                        page_text = "\n".join(page_text_parts).strip()
                        is_scanned = self._detect_scanned_page(page_text, len(images))
                        _ocr_content = getattr(getattr(Config, "VisionOCR", None), "OCR_CONTENT_IMAGES", True)
                        _min_w = getattr(getattr(Config, "VisionOCR", None), "MIN_IMAGE_WIDTH", 100)
                        _min_h = getattr(getattr(Config, "VisionOCR", None), "MIN_IMAGE_HEIGHT", 100)

                        if images and (is_scanned or _ocr_content):
                            for img_index, img in enumerate(images):
                                try:
                                    base_image = doc.extract_image(img[0])
                                    image = Image.open(io.BytesIO(base_image["image"]))
                                    w, h = image.size
                                    if w < _min_w or h < _min_h:
                                        figures.append(Figure(
                                            page=page_index,
                                            caption=f"Image_{page_index}_{img_index}",
                                        ))
                                        continue

                                    ocr_text, ocr_conf = self._ocr_with_vision(
                                        image, is_full_page=is_scanned,
                                    )
                                    ocr_method = "vision_ocr" if getattr(
                                        getattr(Config, "VisionOCR", None), "ENABLED", True,
                                    ) else "pytesseract"
                                    if ocr_text:
                                        chunk_type = "ocr_text" if is_scanned else "image_content"
                                        fig = Figure(
                                            page=page_index,
                                            caption=ocr_text,
                                            ocr_method=ocr_method,
                                            ocr_confidence=ocr_conf,
                                        )
                                        figures.append(fig)
                                        chunk_candidates.append(
                                            ChunkCandidate(
                                                text=ocr_text,
                                                page=page_index,
                                                section_title=current_title,
                                                section_id=current_section_id,
                                                chunk_type=chunk_type,
                                            )
                                        )
                                        # Diagram detection and structural extraction
                                        try:
                                            from src.doc_understanding.diagram_extractor import is_likely_diagram, extract_diagram_structure
                                            from src.api.config import Config as _DiagCfg
                                            if getattr(getattr(_DiagCfg, "DiagramExtraction", None), "ENABLED", True):
                                                if ocr_text and is_likely_diagram(ocr_text, (w, h)):
                                                    use_think = getattr(getattr(_DiagCfg, "DiagramExtraction", None), "USE_THINKING", True)
                                                    diagram_struct = extract_diagram_structure(ocr_text, use_thinking=use_think)
                                                    if diagram_struct:
                                                        fig.is_diagram = True
                                                        fig.diagram_type = diagram_struct.diagram_type
                                                        fig.diagram_structure = diagram_struct.to_dict()
                                                        # Add diagram as a rich chunk candidate
                                                        chunk_candidates.append(ChunkCandidate(
                                                            text=diagram_struct.to_text(),
                                                            page=page_index,
                                                            section_title=f"Diagram: {diagram_struct.diagram_type}",
                                                            section_id=None,
                                                            chunk_type="diagram",
                                                        ))
                                                        logger.info("Extracted %s diagram with %d nodes on page %d",
                                                            diagram_struct.diagram_type, len(diagram_struct.nodes), page_index)
                                        except Exception as _diag_exc:
                                            logger.debug("Diagram extraction skipped: %s", _diag_exc)
                                        page_text_parts.append(ocr_text)
                                        ocr_target = {
                                            "page": page_index,
                                            "section_title": current_title,
                                            "section_id": current_section_id,
                                            "image": image.copy(),
                                            "candidate_index": len(chunk_candidates) - 1,
                                            "figure_index": len(figures) - 1,
                                            "original_text": ocr_text,
                                            "original_conf": ocr_conf,
                                        }
                                        if ocr_conf is not None:
                                            ocr_confidences.append(float(ocr_conf))
                                            ocr_target["conf_index"] = len(ocr_confidences) - 1
                                        ocr_targets.append(ocr_target)
                                    else:
                                        figures.append(Figure(
                                            page=page_index,
                                            caption=f"Image_{page_index}_{img_index}",
                                        ))
                                except Exception as exc:  # noqa: BLE001
                                    logger.debug("OCR image extraction failed on page %s: %s", page_index, exc)
                                    errors.append(f"image_extraction_failed: page={page_index} err={exc}")
                        else:
                            for img_index, _ in enumerate(images):
                                figures.append(Figure(page=page_index, caption=f"Image_{page_index}_{img_index}"))

                        if page_text_parts:
                            pages_map[page_index] = list(page_text_parts)
                            full_text_parts.append(f"\n--- Page {page_index} ---\n" + "\n".join(page_text_parts))
            except Exception as exc:  # noqa: BLE001
                errors.append(f"pdf_layout_parse_failed: {exc}")
        if pdfplumber:
            try:
                with pdfplumber.open(io.BytesIO(pdf_content)) as doc:
                    for page_index, page in enumerate(doc.pages, start=1):
                        last_page = max(last_page, page_index)
                        # Skip pages already handled well by fitz+OCR (>50 chars of content)
                        existing_page = pages_map.get(page_index, [])
                        existing_len = sum(len(line) for line in existing_page) if existing_page else 0
                        text = page.extract_text() or ""
                        if text.strip() and len(text.strip()) > existing_len:
                            page_lines = text.splitlines()
                            for para in text.split("\n"):
                                para = para.strip()
                                if not para:
                                    continue
                                if self._is_heading(para):
                                    self._finalize_section(
                                        sections, current_title, current_section_id, current_start_page, page_index, section_buffer
                                    )
                                    section_buffer = []
                                    current_title = para
                                    current_section_id = self._make_section_id(para, page_index)
                                    current_start_page = page_index
                                    continue
                                section_buffer.append(para)
                                chunk_type = "table" if self._looks_like_table(para) else "text"
                                chunk_candidates.append(
                                    ChunkCandidate(
                                        text=para,
                                        page=page_index,
                                        section_title=current_title,
                                        section_id=current_section_id,
                                        chunk_type=chunk_type,
                                    )
                                )
                                if chunk_type == "table":
                                    tables.append(Table(page=page_index, text=para))
                            full_text_parts.append(f"\n--- Page {page_index} ---\n{text}")
                            pages_map[page_index] = page_lines

                        try:
                            extracted_tables = page.extract_tables()
                            if extracted_tables:
                                try:
                                    from src.intelligence.table_parser import TableParser
                                    _table_parser = TableParser()
                                except Exception:
                                    _table_parser = None
                                for tbl in extracted_tables:
                                    if _table_parser:
                                        structured = _table_parser.parse(tbl, page=page_index, title=current_title)
                                        table_type = _table_parser.classify_table_type(structured)
                                        structured.table_type = table_type
                                        formatted_table = structured.flat_text
                                        tables.append(Table(page=page_index, text=formatted_table, csv=formatted_table, structured=structured))
                                        chunk_candidates.append(
                                            ChunkCandidate(
                                                text=formatted_table,
                                                page=page_index,
                                                section_title=current_title,
                                                section_id=current_section_id,
                                                chunk_type="table",
                                                table_meta={
                                                    "headers": structured.headers,
                                                    "table_type": table_type,
                                                    "row_count": structured.row_count,
                                                    "col_count": structured.col_count,
                                                },
                                            )
                                        )
                                    else:
                                        formatted_table = "\n".join([", ".join(str(c) for c in row) for row in tbl if row])
                                        tables.append(Table(page=page_index, text=formatted_table, csv=formatted_table))
                                        chunk_candidates.append(
                                            ChunkCandidate(
                                                text=formatted_table,
                                                page=page_index,
                                                section_title=current_title,
                                                section_id=current_section_id,
                                                chunk_type="table",
                                            )
                                        )
                        except Exception as exc:  # noqa: BLE001
                            errors.append(f"pdf_table_extract_failed: page={page_index} err={exc}")
            except Exception as exc:  # noqa: BLE001
                errors.append(f"pdfplumber_open_failed: {exc}")
        else:
            errors.append("No PDF parser available: PyMuPDF/pdfplumber missing")

        self._finalize_section(
            sections, current_title, current_section_id, current_start_page, last_page, section_buffer
        )

        if pages_map:
            cleaned_pages, header_lines, footer_lines = self._dedupe_page_lines(pages_map)
            if cleaned_pages:
                pages_map = cleaned_pages
                full_text_parts = [
                    f"\n--- Page {page} ---\n" + "\n".join(lines)
                    for page, lines in sorted(cleaned_pages.items())
                    if lines
                ]
                if header_lines or footer_lines:
                    chunk_candidates = [
                        cand
                        for cand in chunk_candidates
                        if cand.text.strip() not in header_lines
                        and cand.text.strip() not in footer_lines
                        and not self._page_number_line(cand.text.strip())
                    ]

        quality_snapshot = self._assess_doc_quality(
            ocr_confidences=ocr_confidences,
            chunk_candidates=chunk_candidates,
            has_images=bool(ocr_targets),
        )
        if (
            quality_snapshot["doc_quality"] == "LOW"
            and ocr_targets
            and easyocr
            and self.ocr_engine != "easyocr"
        ):
            for target in ocr_targets:
                new_text, new_conf = self._ocr_image(target["image"], engine="easyocr")
                if not new_text:
                    continue
                old_text = target.get("original_text") or ""
                old_conf = target.get("original_conf") or 0.0
                if new_conf is None:
                    new_conf = old_conf
                if len(new_text) <= len(old_text) and new_conf <= old_conf:
                    continue
                idx = target.get("candidate_index")
                if isinstance(idx, int) and 0 <= idx < len(chunk_candidates):
                    chunk_candidates[idx].text = new_text
                fig_idx = target.get("figure_index")
                if isinstance(fig_idx, int) and 0 <= fig_idx < len(figures):
                    figures[fig_idx].caption = new_text
                conf_idx = target.get("conf_index")
                if isinstance(conf_idx, int) and 0 <= conf_idx < len(ocr_confidences):
                    ocr_confidences[conf_idx] = float(new_conf)

            quality_snapshot = self._assess_doc_quality(
                ocr_confidences=ocr_confidences,
                chunk_candidates=chunk_candidates,
                has_images=bool(ocr_targets),
            )

        full_text = "\n".join(full_text_parts).strip()

        # Document-level OCR fallback: if no meaningful text was extracted,
        # force OCR on ALL pages (the PDF is likely fully image-based)
        _text_stripped = full_text.strip()
        _alnum_ratio = sum(c.isalnum() for c in _text_stripped) / max(len(_text_stripped), 1)
        _text_is_garbage = (
            len(_text_stripped) < 100
            or (len(_text_stripped) < 800 and _alnum_ratio < 0.5)
        )
        if _text_is_garbage and not chunk_candidates:
            logger.info(
                "PDF extraction produced low-quality text (%d chars, %.0f%% alnum); triggering full-document OCR",
                len(_text_stripped), _alnum_ratio * 100,
            )
            try:
                if hasattr(fitz, "open"):
                    with fitz.open(stream=pdf_content, filetype="pdf") as doc:
                        for page_index, page in enumerate(doc, start=1):
                            last_page = max(last_page, page_index)
                            # Render page as image for OCR
                            try:
                                mat = fitz.Matrix(2.0, 2.0)  # 2x zoom for better OCR
                                pix = page.get_pixmap(matrix=mat)
                                img_data = pix.tobytes("png")
                                image = Image.open(io.BytesIO(img_data))
                                ocr_text, ocr_conf = self._ocr_with_vision(image, is_full_page=True)
                                if ocr_text and ocr_text.strip():
                                    full_text_parts.append(f"\n--- Page {page_index} ---\n{ocr_text}")
                                    for para in ocr_text.split("\n"):
                                        para = para.strip()
                                        if not para:
                                            continue
                                        section_buffer.append(para)
                                        chunk_type = "table" if self._looks_like_table(para) else "text"
                                        chunk_candidates.append(
                                            ChunkCandidate(
                                                text=para,
                                                page=page_index,
                                                section_title=current_title,
                                                section_id=current_section_id,
                                                chunk_type=chunk_type,
                                            )
                                        )
                                    pages_map[page_index] = ocr_text.splitlines()
                                    if ocr_conf is not None:
                                        ocr_confidences.append(float(ocr_conf))
                            except Exception as exc:  # noqa: BLE001
                                logger.debug("Full-page OCR failed for page %s: %s", page_index, exc)
                                errors.append(f"full_page_ocr_failed: page={page_index} err={exc}")
                    # Finalize remaining section buffer from OCR
                    self._finalize_section(
                        sections, current_title, current_section_id, current_start_page, last_page, section_buffer
                    )
                    section_buffer = []
                    full_text = "\n".join(full_text_parts).strip()
                    logger.info("Full-document OCR produced %d chars, %d chunks", len(full_text), len(chunk_candidates))
            except Exception as exc:  # noqa: BLE001
                logger.warning("Full-document OCR fallback failed: %s", exc)
                errors.append(f"full_doc_ocr_failed: {exc}")

        doc_type = self._doc_intel.infer_type(tables, figures, sections, full_text, filename_hint=filename or "document.pdf")
        pages = [
            {"page_number": page, "text": "\n".join(lines).strip()}
            for page, lines in sorted(pages_map.items())
            if lines
        ]
        canonical_json = self._build_canonical_json(
            pages=pages,
            sections=sections,
            tables=tables,
            figures=figures,
            chunk_candidates=chunk_candidates,
            layout_blocks=layout_blocks,
            page_dims=page_dims,
        )
        return ExtractedDocument(
            full_text=full_text,
            sections=sections,
            tables=tables,
            figures=figures,
            chunk_candidates=chunk_candidates,
            doc_type=doc_type,
            errors=errors,
            metrics={
                "ocr_confidences": ocr_confidences,
                **quality_snapshot,
            },
            canonical_json=canonical_json,
            doc_quality=quality_snapshot.get("doc_quality"),
        )

    def extract_text_from_docx(self, doc_content: Union[bytes, bytearray, memoryview, str, io.IOBase], filename: Optional[str] = None) -> ExtractedDocument:
        document = docx.Document(self._coerce_file_like(doc_content))
        sections: List[Section] = []
        tables: List[Table] = []
        figures: List[Figure] = []
        chunk_candidates: List[ChunkCandidate] = []
        full_text_parts: List[str] = []
        errors: List[str] = []

        current_title = "Introduction"
        current_section_id = self._make_section_id(current_title, 1)
        section_buffer: List[str] = []

        for para in document.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style = para.style.name if para.style else ""
            if style.startswith("Heading") or self._is_heading(text):
                self._finalize_section(sections, current_title, current_section_id, 1, 1, section_buffer)
                section_buffer = []
                current_title = text
                current_section_id = self._make_section_id(current_title, 1)
                continue

            section_buffer.append(text)
            chunk_candidates.append(
                ChunkCandidate(
                    text=text,
                    page=None,
                    section_title=current_title,
                    section_id=current_section_id,
                    chunk_type="text",
                )
            )
            full_text_parts.append(text)

        # Tables
        for tbl in document.tables:
            rows = []
            for row in tbl.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    rows.append(" | ".join(row_text))
            if rows:
                csv_like = "\n".join(rows)
                tables.append(Table(page=None, text=csv_like, csv=csv_like))
                chunk_candidates.append(
                    ChunkCandidate(
                        text=csv_like,
                        page=None,
                        section_title=current_title,
                        section_id=current_section_id,
                        chunk_type="table",
                    )
                )
                full_text_parts.append(csv_like)

        self._finalize_section(sections, current_title, current_section_id, 1, 1, section_buffer)

        # Figures via document relationships (captions/alt-text if present)
        try:
            for rel in document.part.rels.values():
                if "image" in rel.target_ref:
                    caption = getattr(rel, "alt_text", None) or rel.target_ref
                    figures.append(Figure(page=None, caption=caption))
        except Exception as exc:  # noqa: BLE001
            errors.append(f"docx_image_parse_failed: {exc}")

        try:
            for shape in getattr(document, "inline_shapes", []):
                alt_text = getattr(shape, "alt_text", "") or getattr(shape, "description", "")
                if alt_text:
                    figures.append(Figure(page=None, caption=alt_text))
        except Exception as exc:  # noqa: BLE001
            errors.append(f"docx_inline_shapes_failed: {exc}")

        full_text = "\n".join(full_text_parts).strip()
        doc_type = self._doc_intel.infer_type(tables, figures, sections, full_text, filename_hint=filename or "document.docx")
        quality_snapshot = self._assess_doc_quality(
            ocr_confidences=[],
            chunk_candidates=chunk_candidates,
            has_images=bool(figures),
        )
        pages = [{"page_number": 1, "text": full_text}] if full_text else []
        canonical_json = self._build_canonical_json(
            pages=pages,
            sections=sections,
            tables=tables,
            figures=figures,
            chunk_candidates=chunk_candidates,
        )
        return ExtractedDocument(
            full_text=full_text,
            sections=sections,
            tables=tables,
            figures=figures,
            chunk_candidates=chunk_candidates,
            doc_type=doc_type,
            errors=errors,
            metrics=quality_snapshot,
            canonical_json=canonical_json,
            doc_quality=quality_snapshot.get("doc_quality"),
        )

    def extract_text_from_pptx(self, ppt_content: Union[bytes, bytearray, memoryview, str, io.IOBase], filename: Optional[str] = None) -> ExtractedDocument:
        presentation = Presentation(self._coerce_file_like(ppt_content))
        sections: List[Section] = []
        tables: List[Table] = []
        figures: List[Figure] = []
        chunk_candidates: List[ChunkCandidate] = []
        full_text_parts: List[str] = []
        errors: List[str] = []

        for slide_num, slide in enumerate(presentation.slides, start=1):
            slide_title = f"Slide {slide_num}"
            section_id = self._make_section_id(slide_title, slide_num)
            slide_text_parts: List[str] = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    slide_text_parts.append(text)
                    chunk_candidates.append(
                        ChunkCandidate(
                            text=text,
                            page=slide_num,
                            section_title=slide_title,
                            section_id=section_id,
                            chunk_type="text",
                        )
                    )

                if shape.has_table:
                    rows = []
                    for row in shape.table.rows:
                        row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_text:
                            rows.append(" | ".join(row_text))
                    if rows:
                        csv_like = "\n".join(rows)
                        tables.append(Table(page=slide_num, text=csv_like, csv=csv_like))
                        chunk_candidates.append(
                            ChunkCandidate(
                                text=csv_like,
                                page=slide_num,
                                section_title=slide_title,
                                section_id=section_id,
                                chunk_type="table",
                            )
                        )
                        slide_text_parts.append(csv_like)

                if getattr(shape, "shape_type", None) == 13 and getattr(shape, "image", None):
                    caption = getattr(shape, "name", "") or shape.image.filename or "Figure"
                    figures.append(Figure(page=slide_num, caption=caption))

            if slide_text_parts:
                full_text_parts.append(f"\n--- Slide {slide_num} ---\n" + "\n".join(slide_text_parts))
                sections.append(
                    Section(
                        section_id=section_id,
                        title=slide_title,
                        level=1,
                        start_page=slide_num,
                        end_page=slide_num,
                        text="\n".join(slide_text_parts),
                    )
                )

        full_text = "\n".join(full_text_parts).strip()
        doc_type = self._doc_intel.infer_type(tables, figures, sections, full_text, filename_hint=filename or "slides.pptx")
        quality_snapshot = self._assess_doc_quality(
            ocr_confidences=[],
            chunk_candidates=chunk_candidates,
            has_images=bool(figures),
        )
        pages = [
            {"page_number": sec.start_page or 1, "text": sec.text}
            for sec in sections
            if sec.text
        ]
        canonical_json = self._build_canonical_json(
            pages=pages,
            sections=sections,
            tables=tables,
            figures=figures,
            chunk_candidates=chunk_candidates,
        )
        return ExtractedDocument(
            full_text=full_text,
            sections=sections,
            tables=tables,
            figures=figures,
            chunk_candidates=chunk_candidates,
            doc_type=doc_type,
            errors=errors,
            metrics=quality_snapshot,
            canonical_json=canonical_json,
            doc_quality=quality_snapshot.get("doc_quality"),
        )

    def extract_text_from_txt(self, text_content: Union[bytes, str], filename: Optional[str] = None) -> ExtractedDocument:
        """
        Normalize plain text into the structured ExtractedDocument used by other parsers.
        This keeps chunking/metadata consistent with richer document types.
        """
        errors: List[str] = []
        try:
            if isinstance(text_content, (bytes, bytearray)):
                text = _smart_decode(text_content)
            else:
                text = str(text_content)
        except Exception as exc:  # noqa: BLE001
            errors.append(f"txt_decode_failed: {exc}")
            text = ""

        text = text.replace("\r\n", "\n").strip()
        if not text:
            canonical_json = self._build_canonical_json(
                pages=[],
                sections=[],
                tables=[],
                figures=[],
                chunk_candidates=[],
            )
            return ExtractedDocument(
                full_text="",
                sections=[],
                tables=[],
                figures=[],
                chunk_candidates=[],
                doc_type="document",
                errors=errors or ["empty_text"],
                metrics={"doc_quality": "LOW"},
                canonical_json=canonical_json,
                doc_quality="LOW",
            )

        paragraphs = [para.strip() for para in re.split(r"\n\s*\n", text) if para.strip()]
        section_title = "Text Document"
        section_id = self._make_section_id(section_title, 1)

        chunk_candidates: List[ChunkCandidate] = []
        for para in paragraphs:
            chunk_candidates.append(
                ChunkCandidate(
                    text=para,
                    page=None,
                    section_title=section_title,
                    section_id=section_id,
                    chunk_type="text",
                )
            )

        sections = [
            Section(
                section_id=section_id,
                title=section_title,
                level=1,
                start_page=1,
                end_page=1,
                text="\n\n".join(paragraphs),
            )
        ]

        full_text = "\n\n".join(paragraphs)
        doc_type = self._doc_intel.infer_type([], [], sections, full_text, filename_hint=filename or "document.txt")
        quality_snapshot = self._assess_doc_quality(
            ocr_confidences=[],
            chunk_candidates=chunk_candidates,
            has_images=False,
        )
        pages = [{"page_number": 1, "text": full_text}] if full_text else []
        canonical_json = self._build_canonical_json(
            pages=pages,
            sections=sections,
            tables=[],
            figures=[],
            chunk_candidates=chunk_candidates,
        )
        return ExtractedDocument(
            full_text=full_text,
            sections=sections,
            tables=[],
            figures=[],
            chunk_candidates=chunk_candidates,
            doc_type=doc_type,
            errors=errors,
            metrics=quality_snapshot,
            canonical_json=canonical_json,
            doc_quality=quality_snapshot.get("doc_quality"),
        )

    def extract_dataframe(self, df: pd.DataFrame, sheet_name: str = "Sheet1", filename: Optional[str] = None) -> ExtractedDocument:
        """
        Build a structured ExtractedDocument from tabular data (CSV/Excel).
        Preserves headers, row-level detail, and a CSV rendition for downstream chunking.
        """
        try:
            df = df.copy()
            df = df.fillna("")
            df.columns = df.columns.astype(str)
        except Exception as exc:  # noqa: BLE001
            logger.error("Failed to normalize dataframe for extraction: %s", exc)
            raise

        # Create a CSV-like text block for reference and table capture
        table_csv = df.to_csv(index=False)
        tables = [Table(page=None, text=table_csv, csv=table_csv)]

        sections: List[Section] = []
        chunk_candidates: List[ChunkCandidate] = []
        figures: List[Figure] = []
        full_text_parts: List[str] = []

        # Build a simple section representing this sheet/table
        section_id = self._make_section_id(sheet_name, 1)

        # Generate row-wise text for richer chunking
        row_texts: List[str] = []
        headers = list(df.columns)
        for idx, row in df.iterrows():
            pairs = [f"{col}: {row[col]}" for col in headers if str(row[col]).strip() != ""]
            row_text = "; ".join(pairs).strip()
            if not row_text:
                continue
            row_texts.append(row_text)
            chunk_candidates.append(
                ChunkCandidate(
                    text=row_text,
                    page=None,
                    section_title=sheet_name,
                    section_id=section_id,
                    chunk_type="table_row",
                )
            )

        # Add column-summary chunk to retain header context
        header_text = " | ".join(headers)
        if header_text:
            chunk_candidates.insert(
                0,
                ChunkCandidate(
                    text=f"Columns: {header_text}",
                    page=None,
                    section_title=sheet_name,
                    section_id=section_id,
                    chunk_type="table_header",
                ),
            )
            full_text_parts.append(f"Columns: {header_text}")

        if row_texts:
            full_text_parts.append("\n".join(row_texts))

        sections.append(
            Section(
                section_id=section_id,
                title=sheet_name,
                level=1,
                start_page=1,
                end_page=1,
                text="\n".join(row_texts) if row_texts else header_text,
            )
        )

        full_text_parts.append(table_csv)
        full_text = "\n".join([part for part in full_text_parts if part]).strip()

        doc_type = self._doc_intel.infer_type(tables, figures, sections, full_text, filename_hint=filename or sheet_name)
        quality_snapshot = self._assess_doc_quality(
            ocr_confidences=[],
            chunk_candidates=chunk_candidates,
            has_images=False,
        )
        pages = [{"page_number": 1, "text": full_text}] if full_text else []
        canonical_json = self._build_canonical_json(
            pages=pages,
            sections=sections,
            tables=tables,
            figures=figures,
            chunk_candidates=chunk_candidates,
        )
        return ExtractedDocument(
            full_text=full_text,
            sections=sections,
            tables=tables,
            figures=figures,
            chunk_candidates=chunk_candidates,
            doc_type=doc_type,
            errors=[],
            metrics=quality_snapshot,
            canonical_json=canonical_json,
            doc_quality=quality_snapshot.get("doc_quality"),
        )
class DocumentIntelligence:
    """Lightweight classifier to tag document type from observed layout/content."""

    @staticmethod
    def infer_type(
        tables: List[Table],
        figures: List[Figure],
        sections: List[Section],
        full_text: str,
        filename_hint: Optional[str] = None,
    ) -> str:
        # Strategy 1: filename + structure heuristic
        name = (filename_hint or "").lower()
        if name.endswith((".ppt", ".pptx")):
            return "presentation"
        if name.endswith(".docx"):
            return "document"
        if name.endswith(".pdf"):
            # Heuristics based on structure
            if len(tables) >= 3:
                return "report_table_heavy"
            if len(figures) >= 3:
                return "report_visual"
            if sections and len(sections) <= 5 and any("slide" in s.title.lower() for s in sections):
                return "presentation_pdf"
            return "report"
        if tables and not figures:
            return "table_dataset"
        if figures and not tables:
            return "visual_brief"
        # Fallback to narrative doc
        return "document"
