
import io
import fitz
import docx
import logging
import pytesseract
import numpy as np
import pandas as pd
from PIL import Image
from pptx import Presentation
from sklearn.preprocessing import MinMaxScaler

logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")

pytesseract.pytesseract.tesseract_cmd = "/usr/bin/tesseract"

class DocumentExtractor(object):

    def extract_text_from_pdf(self, pdf_content):
        """Extracts text from a PDF file, including OCR for images."""
        try:
            logging.info("Extracting text from PDF (OCR applied only if images exist)...")
            extracted_text = []

            with fitz.open(stream=pdf_content, filetype="pdf") as doc:
                logging.info("Analyzing PDF contents")
                for page_num, page in enumerate(doc, start=1):
                    text = page.get_text("text").strip()
                    extracted_text.append(f"\n--- Page {page_num} ---\n{text}" if text else "")
                    logging.info(f"Reading page {page_num}")

                    images = page.get_images(full=True)
                    if images:
                        logging.info(f"Images found on page {page_num}: {len(images)}")

                        for img_index, img in enumerate(images[:5]):
                            logging.info(f"Extracting image {img_index + 1} from page {page_num}")
                            try:
                                xref = img[0]
                                base_image = doc.extract_image(xref)
                                image_bytes = base_image["image"]
                                image = Image.open(io.BytesIO(image_bytes))

                                image = image.convert("L")
                                image = image.resize((image.width // 2, image.height // 2))

                                ocr_text = pytesseract.image_to_string(image).strip()
                                if ocr_text:
                                    extracted_text.append(f"\n[OCR Extracted from Image {img_index + 1}]\n{ocr_text}")

                            except Exception as img_error:
                                logging.error(f"Error processing image {img_index + 1}: {img_error}")

            result_text = "\n".join(filter(None, extracted_text)).strip()
            return result_text if result_text else "No text found in the PDF."

        except Exception as e:
            logging.error(f"Failed to extract text from PDF: {e}")
            return ""


    def extract_text_from_docx(self, doc_content):
        """Extracts text from a DOCX file, including OCR for images."""
        try:
            logging.info("Extracting text from DOCX (OCR applied only if images exist)...")
            doc = docx.Document(doc_content)
            extracted_text = []

            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    extracted_text.append(text)

            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        extracted_text.append(" | ".join(row_text))

            for rel in doc.part.rels:
                if "image" in doc.part.rels[rel].target_ref:
                    image_data = doc.part.rels[rel].target_part.blob
                    try:
                        image = Image.open(io.BytesIO(image_data))
                        ocr_text = pytesseract.image_to_string(image).strip()
                        if ocr_text:
                            extracted_text.append(f"\n[OCR Extracted from Image]\n{ocr_text}")
                    except Exception as img_error:
                        logging.error(f"Error processing image: {img_error}")

            result_text = "\n".join(filter(None, extracted_text)).strip()

            return result_text if result_text else "No text found in the DOCX."

        except Exception as e:
            logging.error(f"Failed to extract text from DOCX: {e}")
            return ""


    def extract_text_from_pptx(self, ppt_content):
        """Extracts structured text from a PowerPoint (PPTX) file, including OCR for images."""
        try:
            logging.info("Extracting text from PPTX (OCR applied only if images exist)...")
            presentation = Presentation(ppt_content)
            extracted_text = []

            for slide_num, slide in enumerate(presentation.slides, start=1):
                slide_text = [f"\n--- Slide {slide_num} ---"]

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if row_text:
                                slide_text.append(" | ".join(row_text))

                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Type 13 is a Picture
                        image_stream = io.BytesIO(shape.image.blob)
                        image = Image.open(image_stream)
                        ocr_text = pytesseract.image_to_string(image).strip()
                        if ocr_text:
                            slide_text.append(f"\n[OCR Extracted from Image]\n{ocr_text}")

                if len(slide_text) > 1:
                    extracted_text.append("\n".join(slide_text))

            result_text = "\n".join(filter(None, extracted_text)).strip()

            return result_text if result_text else "No text found in the PPTX."

        except Exception as e:
            logging.error(f"Failed to extract text from PPTX: {e}")
            return ""

    def extract_dataframe(self, df, MODEL):
        try:
            df = df.select_dtypes(include=[np.number, 'object']).copy()
            for col in df.select_dtypes(include=['object']).columns:
                try:
                    df[col] = pd.to_datetime(df[col])
                    df[f"Days Since {col}"] = (df[col] - df[col].min()).dt.days
                    df.drop(columns=[col], inplace=True)
                except Exception:
                    logging.warning(f"Column {col} could not be converted to datetime. Keeping it as categorical.")
                    df[col] = df[col].astype(str)

            df.columns = df.columns.astype(str)
            categorical_cols = df.select_dtypes(include=['object']).columns.tolist()
            if categorical_cols:
                text_data = df[categorical_cols].astype(str).apply(lambda x: ' '.join(x.dropna()), axis=1).tolist()
                embeddings = MODEL.encode(text_data, convert_to_numpy=True)
                return {"embeddings": embeddings, "texts": text_data}

        except Exception as e:
            logging.error(f"Error processing dataframe: {e}")
            return None

#final code